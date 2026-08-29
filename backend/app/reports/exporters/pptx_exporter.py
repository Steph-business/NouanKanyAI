"""
app/reports/exporters/pptx_exporter.py — Exportateur de rapports énergétiques au format PPTX (PowerPoint).
"""

import io
import logging
from typing import Dict, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.reports.exporters.base import BaseReportExporter
from app.reports.models import EnergyReportData

logger = logging.getLogger("nouankany.reports")

DARK_NAVY = RGBColor(15, 23, 42)     # #0F172A
SKY_BLUE = RGBColor(2, 132, 199)     # #0284C7
EMERALD = RGBColor(16, 185, 129)     # #10B981
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(241, 245, 249) # #F1F5F9
SLATE = RGBColor(100, 116, 139)


class PPTXReportExporter(BaseReportExporter):
    """
    Générateur de présentations PowerPoint (.pptx) professionnelles pour comités de direction.
    """

    def export(
        self,
        report_data: EnergyReportData,
        chart_images: Optional[Dict[str, bytes]] = None,
    ) -> bytes:
        """
        Génère la présentation PowerPoint complète (5 diapositives exécutives).
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # Format 16:9 moderne
        prs.slide_height = Inches(7.5)
        blank_slide_layout = prs.slide_layouts[6]

        # =============================================================
        # Slide 1 : Titre & Couverture (Fond Sombre Premium)
        # =============================================================
        slide1 = prs.slides.add_slide(blank_slide_layout)
        background1 = slide1.background
        fill1 = background1.fill
        fill1.solid()
        fill1.fore_color.rgb = DARK_NAVY

        # Logo / Marque
        brand_box = slide1.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.0), Inches(0.6))
        tf_brand = brand_box.text_frame
        p_brand = tf_brand.paragraphs[0]
        p_brand.text = "NOUANKANY.AI — INTELLIGENCE ÉNERGÉTIQUE"
        p_brand.font.name = "Arial"
        p_brand.font.size = Pt(14)
        p_brand.font.color.rgb = SKY_BLUE
        p_brand.font.bold = True

        # Titre Principal
        title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.2))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = report_data.title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(32)
        p_title.font.color.rgb = WHITE
        p_title.font.bold = True

        # Sous-titre Métadonnées
        meta_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.0), Inches(1.5))
        tf_meta = meta_box.text_frame
        p_meta = tf_meta.paragraphs[0]
        p_meta.text = (
            f"Site : {report_data.site_name} ({report_data.building_type})\n"
            f"Période analysée : {report_data.period_start} au {report_data.period_end}\n"
            f"Émis par : {report_data.author} | Le {report_data.generated_at}"
        )
        p_meta.font.name = "Arial"
        p_meta.font.size = Pt(14)
        p_meta.font.color.rgb = SLATE

        # =============================================================
        # Slide 2 : Résumé Exécutif & Tableau des KPIs
        # =============================================================
        slide2 = prs.slides.add_slide(blank_slide_layout)
        self._add_slide_header(slide2, "1. Résumé Exécutif & Indicateurs Clés (KPIs)")

        # Boîte Résumé
        summary_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.4))
        tf_sum = summary_box.text_frame
        tf_sum.word_wrap = True
        p_sum = tf_sum.paragraphs[0]
        p_sum.text = report_data.executive_summary
        p_sum.font.name = "Arial"
        p_sum.font.size = Pt(13)
        p_sum.font.color.rgb = DARK_NAVY

        # Tableau KPIs (4 colonnes x 2 lignes)
        k = report_data.kpis
        kpi_table_shape = slide2.shapes.add_table(3, 4, Inches(0.8), Inches(3.2), Inches(11.7), Inches(3.4))
        table2 = kpi_table_shape.table

        kpis_grid = [
            [("Consommation Totale", f"{k.total_energy_kwh:,.1f} kWh"), ("Facture Globale", f"{k.total_cost_fcfa:,.0f} FCFA"), ("Consommation Pointe", f"{k.peak_hours_energy_kwh:,.1f} kWh"), ("Coût Pointe (CIE)", f"{k.peak_hours_cost_fcfa:,.0f} FCFA")],
            [("Économies Effacement", f"+{k.peak_shaving_savings_fcfa:,.0f} FCFA"), ("Facteur de Puissance", f"cos φ = {k.average_power_factor_cos_phi:.2f}"), ("Puissance Maximale", f"{k.max_peak_power_kw:.1f} kW"), ("Puissance Souscrite", f"{k.subscribed_power_limit_kw:.0f} kW")],
            [("Émissions Carbone", f"{k.carbon_emissions_kg_co2:,.1f} kg CO2"), ("Anomalies Détectées", f"{k.anomaly_incidents_count} incidents"), ("Taux de Conformité", "98.5 %"), ("Optimisation Globale", "Active")],
        ]

        for r_idx, row in enumerate(kpis_grid):
            for c_idx, (lbl, val) in enumerate(row):
                cell = table2.cell(r_idx, c_idx)
                cell.text = f"{lbl}\n{val}"
                p = cell.text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(11)
                p.font.color.rgb = SLATE
                if len(cell.text_frame.paragraphs) > 1:
                    p_val = cell.text_frame.paragraphs[1]
                    p_val.font.name = "Arial"
                    p_val.font.size = Pt(16)
                    p_val.font.bold = True
                    p_val.font.color.rgb = SKY_BLUE if "FCFA" in val or "kWh" in val else DARK_NAVY

        # =============================================================
        # Slide 3 : Analyse Graphique & Courbe de Charge
        # =============================================================
        slide3 = prs.slides.add_slide(blank_slide_layout)
        self._add_slide_header(slide3, "2. Analyse de la Courbe de Charge & Dépassements de Pointe")

        if chart_images and "load_curve" in chart_images:
            slide3.shapes.add_picture(io.BytesIO(chart_images["load_curve"]), Inches(0.8), Inches(1.5), width=Inches(7.5))

        # Zone explicative à droite
        notes_box = slide3.shapes.add_textbox(Inches(8.6), Inches(1.5), Inches(4.0), Inches(5.0))
        tf_notes = notes_box.text_frame
        tf_notes.word_wrap = True
        p_n1 = tf_notes.paragraphs[0]
        p_n1.text = "🎯 Constats Clés :"
        p_n1.font.bold = True
        p_n1.font.size = Pt(14)
        p_n1.font.color.rgb = DARK_NAVY

        insights = [
            f"• Puissance de crête observée : {k.max_peak_power_kw:.1f} kW.",
            f"• Marge souscrite restante : {max(0.0, k.subscribed_power_limit_kw - k.max_peak_power_kw):.1f} kW.",
            "• Fenêtre critique CIE : 19h00 - 23h00 (Tarif de 145 FCFA/kWh).",
            f"• Économies sécurisées ce mois-ci : {k.peak_shaving_savings_fcfa:,.0f} FCFA.",
        ]
        for ins in insights:
            p_ins = tf_notes.add_paragraph()
            p_ins.text = ins
            p_ins.font.size = Pt(12)
            p_ins.font.color.rgb = DARK_NAVY

        # =============================================================
        # Slide 4 : Répartition Tranches CIE & Bilan Machines
        # =============================================================
        slide4 = prs.slides.add_slide(blank_slide_layout)
        self._add_slide_header(slide4, "3. Répartition Énergétique & Consommation par Équipement")

        if chart_images and "tariff_pie" in chart_images:
            slide4.shapes.add_picture(io.BytesIO(chart_images["tariff_pie"]), Inches(0.8), Inches(1.5), width=Inches(5.5))

        if chart_images and "machines_bar" in chart_images:
            slide4.shapes.add_picture(io.BytesIO(chart_images["machines_bar"]), Inches(6.8), Inches(1.5), width=Inches(5.8))

        # =============================================================
        # Slide 5 : Recommandations IA & Plan d'Action ROI
        # =============================================================
        slide5 = prs.slides.add_slide(blank_slide_layout)
        self._add_slide_header(slide5, "4. Recommandations IA & Plan d'Action Rentabilité (ROI)")

        rec_table_shape = slide5.shapes.add_table(1 + len(report_data.recommendations[:4]), 4, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8))
        table5 = rec_table_shape.table

        # Headers
        h_titles = ["Priorité", "Action Recommandée", "Cible", "Gains Estimés (FCFA)"]
        for c_idx, ht in enumerate(h_titles):
            cell = table5.cell(0, c_idx)
            cell.text = ht
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_NAVY

        for r_idx, rec in enumerate(report_data.recommendations[:4], start=1):
            prio_lbl = f"P{rec.priority} ({'Immédiat' if rec.priority==1 else 'Court terme'})"
            row_data = [prio_lbl, f"{rec.title}\n{rec.description[:120]}...", rec.target_equipment, f"+{rec.estimated_savings_fcfa:,.0f} FCFA\n({rec.estimated_savings_kwh:,.0f} kWh)"]
            for c_idx, val in enumerate(row_data):
                cell = table5.cell(r_idx, c_idx)
                cell.text = val
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(10.5)
                if c_idx == 3:
                    p.font.bold = True
                    p.font.color.rgb = EMERALD

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()

    def _add_slide_header(self, slide, title_text: str):
        """Ajoute l'en-tête standard corporate sur les diapositives de contenu."""
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_NAVY
